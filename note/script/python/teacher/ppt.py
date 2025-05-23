from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

'''
本周工作总结
读取1.pptx为模板，设置里面的内容
'''

yue = 5
zhou = 4 # 当前第几周，最大5周
file_name = f"2025年{yue}月第{zhou}周工作总结卜凯凯.pptx"
# title = f"{yue}月第{zhou}周工作总结\n6月第1周工作计划"
title = f"{yue}月第{zhou}周工作总结\n{yue}月第{zhou+1}周工作计划"
stitle = f"25年{yue}月第{zhou}周工作完成情况："
do_work = [
"讲解linux的内核的更换，DNS的配置，演示宝塔面板安装及基本功能",
"讲解大数据技术基础中的数据思维与安全",
"讲解人工智能导论中语音识别、分析与语音克隆",
"批改毕业考试与补考学生的试卷，录入成绩",
"参加学院召开的日常会，备课"]
plan_work = [
"讲解linux操作系统的6种启动模式及忘记ROOT密码如何修改，开机GRUB设置，讲解大数据技术基础中的城市空气质量大数据分析实战",
"讲解linux操作系统的6种启动模式及忘记ROOT密码如何修改，开机GRUB设置，讲解人工智能导论中图片生成视频、文字生成视频",
"讲解linux操作系统的6种启动模式及忘记ROOT密码如何修改，开机GRUB设置",
"讲解linux操作系统的6种启动模式及忘记ROOT密码如何修改，开机GRUB设置",
"讲解大数据技术基础中城市空气质量大数据分析实战，讲解人工智能导论中图片生成视频、文字生成视频"]


if os.path.exists(file_name):
    os.remove(file_name)
# 读取 PPT 文件
prs = Presentation('1.pptx')
# prs.slides[0].shapes[0].text_frame.text="1234"
print(prs.slides[0].shapes[0].text_frame.text)
print(prs.slides[3].shapes[1].text_frame.text)
print(prs.slides[3].shapes[0].table.cell(1, 0).text)

# 获取第一张幻灯片的第一个形状
shape1 = prs.slides[0].shapes[0]
text_frame = shape1.text_frame
text_frame.clear()  # 清空原有内容，但保留段落属性
para = text_frame.paragraphs[0]
# 添加新的文本，并设置字体大小和颜色
new_run = para.add_run()
new_run.text = title
# 设置字体大小和颜色
new_run.font.bold = True
new_run.font.size = Pt(54)  # 设置字体大小为24磅
new_run.font.color.rgb = RGBColor(255, 255, 255)  # 白色
new_run.font.name = '宋体'  


# 获取第三张幻灯片
slide = prs.slides[3]
slide.shapes[1].text_frame.clear()
new_run = slide.shapes[1].text_frame.paragraphs[0].add_run()
new_run.text = stitle
# 设置字体大小和颜色
new_run.font.bold = True
new_run.font.size = Pt(36)
new_run.font.name = '宋体'  

table = prs.slides[3].shapes[0].table
data_list = list(do_work)
for i, value in enumerate(data_list):
        cell_row_index = i + 1  # 从第2行开始写入（索引为1）
        cell = table.cell(cell_row_index, 0)  # 第0列
        cell.text_frame.clear()
        paragraph = cell.text_frame.paragraphs[0]  # clear 后默认保留一个段落
        new_run = paragraph.add_run()
        new_run.text = str(value)  # 转为字符串写入
        new_run.font.size = Pt(14)
        new_run.font.name = '宋体'

# 获取第六张幻灯片的第二个形状
table = prs.slides[5].shapes[1].table
data_list = list(plan_work)
for i, value in enumerate(data_list):
        cell_row_index = i + 1  # 从第2行开始写入（索引为1）
        cell = table.cell(cell_row_index, 1)  # 第0列
        cell.text_frame.clear()
        paragraph = cell.text_frame.paragraphs[0]  # clear 后默认保留一个段落
        new_run = paragraph.add_run()
        new_run.text = str(value)  # 转为字符串写入
        new_run.font.size = Pt(14)
        new_run.font.name = '宋体'

all_shapes = prs.slides[5].shapes

# 遍历所有形状并输出基本信息
for idx, shape in enumerate(all_shapes):
    print(f"形状 {idx + 1}:")
    print(f"  类型: {shape.shape_type}")
    print(f"  名称: {shape.name}")
    print(f"  是否文本框: {shape.has_text_frame}")
    print("-" * 30)

prs.save(file_name)